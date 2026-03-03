from typing import List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from core.document_creator.assemblers.base import BaseDocumentAssembler
from core.document_creator.state import DocumentCreatorConfig, SectionState


class PptxAssembler(BaseDocumentAssembler):
    """
    Assembles sections into a PowerPoint presentation using python-pptx.

    Creates a clean, professional presentation with:
    - Title slide
    - Content slides (one per section)
    - Speaker notes when available
    """

    # Brand colors
    TITLE_COLOR = RGBColor(0x1E, 0x29, 0x3B)  # Dark navy
    SUBTITLE_COLOR = RGBColor(0x64, 0x74, 0x8B)  # Slate gray
    ACCENT_COLOR = RGBColor(0xF9, 0x73, 0x16)  # Orange accent

    async def assemble(
        self,
        title: str,
        subtitle: Optional[str],
        sections: List[SectionState],
        config: DocumentCreatorConfig,
        output_path: str,
    ) -> str:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title slide
        self._add_title_slide(prs, title, subtitle)

        # Content slides
        for section_state in sections:
            if not section_state.versions:
                continue
            version = section_state.versions[section_state.selected_version_index]
            self._add_content_slide(prs, section_state, version)

        prs.save(output_path)
        return output_path

    def _add_title_slide(self, prs, title, subtitle):
        slide_layout = prs.slide_layouts[0]  # Title Slide
        slide = prs.slides.add_slide(slide_layout)

        # Title
        if slide.placeholders[0]:
            tf = slide.placeholders[0]
            tf.text = title
            for paragraph in tf.text_frame.paragraphs:
                paragraph.font.size = Pt(36)
                paragraph.font.color.rgb = self.TITLE_COLOR
                paragraph.font.bold = True

        # Subtitle
        if subtitle and len(slide.placeholders) > 1:
            tf = slide.placeholders[1]
            tf.text = subtitle
            for paragraph in tf.text_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.font.color.rgb = self.SUBTITLE_COLOR

    def _add_content_slide(self, prs, section_state, version):
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)

        # Title
        if slide.placeholders[0]:
            slide.placeholders[0].text = version.heading or section_state.spec.title
            for p in slide.placeholders[0].text_frame.paragraphs:
                p.font.size = Pt(28)
                p.font.color.rgb = self.TITLE_COLOR
                p.font.bold = True

        # Content body
        if len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.clear()

            if version.bullet_points:
                for idx, bullet in enumerate(version.bullet_points):
                    p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                    p.text = bullet
                    p.level = 0
                    p.font.size = Pt(16)
                    p.font.color.rgb = self.TITLE_COLOR
                    p.space_after = Pt(6)
            elif version.content:
                paragraphs = version.content.split("\n")
                for idx, para_text in enumerate(paragraphs):
                    if not para_text.strip():
                        continue
                    p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                    # Truncate very long paragraphs for slides
                    text = para_text.strip()
                    if len(text) > 300:
                        text = text[:297] + "..."
                    p.text = text
                    p.font.size = Pt(14)
                    p.font.color.rgb = self.TITLE_COLOR
                    p.space_after = Pt(4)

        # Speaker notes
        if version.speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = version.speaker_notes

        # Key takeaway as a footer-style note
        if version.key_takeaway:
            from pptx.util import Emu

            left = Inches(0.5)
            top = Inches(6.5)
            width = Inches(12.333)
            height = Inches(0.5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = f"Key Takeaway: {version.key_takeaway}"
            p.font.size = Pt(11)
            p.font.italic = True
            p.font.color.rgb = self.ACCENT_COLOR
